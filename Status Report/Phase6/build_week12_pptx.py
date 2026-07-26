"""Build Week 12 Status Report PPTX — appends previous weeks after Week 12."""

import os, re, shutil, zipfile, tempfile
from pptx import Presentation
from pptx.util import Pt, Emu
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from lxml import etree

TEMPLATE = "/home/ds/MonU_Applied_DataS_cience_Project/Status Report/CS703 - Weekly Status Report Template.pptx"
PREVIOUS = "/home/ds/MonU_Applied_DataS_cience_Project/Status Report/Phase6/CS703 - Week11 Status Report.pptx.submitted"
OUTPUT = "/home/ds/MonU_Applied_DataS_cience_Project/Status Report/Phase6/CS703 - Week12 Status Report.pptx"
SCREENSHOT = "/home/ds/MonU_Applied_DataS_cience_Project/Status Report/Phase6/presentation_slide1.png"

WEEK_ENDING = "July 24, 2026"
SUBMIT_DATE = "July 26, 2026"
SUBMITTED_BY = "Gordon Oboh"
GREEN_HEX = "00FF00"

AGENDA_ITEMS = [
    "Status Overview", "Items Completed This Week", "Items In Progress",
    "Items To Be Started", "Samples of Items Completed This Week",
    "Issues, Risks, Concerns", "Next Steps", "Personal Reflection",
]

REFLECTION = (
    "This is it \u2014 the final status report. Submitting the Final Presentation this week "
    "felt like crossing a major milestone. Seeing the full slide deck come together with "
    "all the required sections made me realize how much ground this project has covered, "
    "from the initial business problem to a working Random Forest model and a live "
    "demonstration. The presentation forced me to distill everything into a clear "
    "narrative, which was challenging but rewarding. I am starting the Final Report now, "
    "and while there is still work ahead, I feel confident because all the pieces are "
    "already written across the Phase 1\u20136 deliverables. It is mostly a matter of "
    "assembling and polishing. This course has been intense, but I have learned a lot "
    "about applying data science to real estate problems and about managing a project "
    "from start to finish."
)

NS = {
    "a": "http://schemas.openxmlformats.org/drawingml/2006/main",
    "r": "http://schemas.openxmlformats.org/officeDocument/2006/relationships",
    "p": "http://schemas.openxmlformats.org/presentationml/2006/main",
}

def seg(text, color=None): return (text, color)
def sts(text): return (text, None)

# ---- Data blocks ----
STATUS_OVERVIEW_LINES = [
    ([sts("Project Status: "), seg("GREEN", GREEN_HEX)], 12, True, 0),
    ("Task 6.3: Final Presentation", 11, True, 0),
    ([sts("Status: "), seg("GREEN", GREEN_HEX), sts(" \u2014 100% complete")], 10, False, 1),
    ("Original Start: 12 Jul 2026 | Revised: N/A | Actual Start: 11 Jul 2026", 9, False, 1),
    ("Original End:   25 Jul 2026 | Revised: N/A | Actual End: 26 Jul 2026", 9, False, 1),
    ("Task 6.3: Final Report", 11, True, 0),
    ([sts("Status: "), seg("GREEN", GREEN_HEX), sts(" \u2014 5% complete")], 10, False, 1),
    ("Original Start: 25 Jul 2026 | Revised: N/A | Actual Start: 26 Jul 2026", 9, False, 1),
    ("Original End:   1 Aug 2026 | Revised: N/A | Actual End: N/A", 9, False, 1),
]

COMPLETED_LINES = [
    ("Task 6.3: Final Presentation", 11, True, 0),
    ([sts("Status: "), seg("GREEN", GREEN_HEX), sts(" \u2014 100% complete")], 10, False, 1),
    ("Original Start: 12 Jul 2026 | Revised: N/A | Actual Start: 11 Jul 2026", 9, False, 1),
    ("Original End:   25 Jul 2026 | Revised: N/A | Actual End: 26 Jul 2026", 9, False, 1),
    ("Final Presentation deliverable completed and submitted.", 10, False, 1),
    ("Slide deck covers all required sections per professor\u2019s outline.", 10, False, 1),
]

IN_PROGRESS_LINES = [
    ("Task 6.3: Final Report", 11, True, 0),
    ([sts("Status: "), seg("GREEN", GREEN_HEX), sts(" \u2014 5% complete")], 10, False, 1),
    ("Original Start: 25 Jul 2026 | Revised: N/A | Actual Start: 26 Jul 2026", 9, False, 1),
    ("Original End:   1 Aug 2026 | Revised: N/A | Actual End: N/A", 9, False, 1),
    ("Initial work begun; outline and section headers drafted,", 10, False, 1),
    ("building on all completed Phase 1\u20136 reports.", 10, False, 1),
]

TBS_LINES = [
    ("Per plan, no new items to start this week.", 11, False, 0),
]

IRC_LINES = [
    ("Issues", 12, True, 0), ("None at this time", 10, False, 1),
    ("", 8, False, 0),
    ("Risks", 12, True, 0), ("None at this time", 10, False, 1),
    ("", 8, False, 0),
    ("Concerns", 12, True, 0), ("None at this time", 10, False, 1),
]

NEXT_STEPS_LINES = [
    ("Task 6.3: Final Report (continue)", 11, True, 0),
    ("Original Start: 25 Jul 2026 | Revised: N/A | Actual Start: 26 Jul 2026", 9, False, 1),
    ("Original End:   1 Aug 2026 | Revised: N/A | Actual End: N/A", 9, False, 1),
]


def _add_paragraph_runs(paragraph, segments, font_size_pt, bold):
    if isinstance(segments, str):
        r = paragraph.add_run()
        r.text = segments; r.font.size = Pt(font_size_pt); r.font.bold = bold
        r.font.name = "Calibri"
        return
    for txt, color_hex in segments:
        r = paragraph.add_run()
        r.text = txt; r.font.size = Pt(font_size_pt); r.font.bold = bold; r.font.name = "Calibri"
        if color_hex is not None:
            r.font.color.rgb = RGBColor(*bytes.fromhex(color_hex))


def fill_content_slide(slide, title, lines):
    for shape in slide.shapes:
        if not shape.has_text_frame: continue
        if shape.placeholder_format and shape.placeholder_format.idx == 0:
            for para in shape.text_frame.paragraphs:
                for run in para.runs: run.text = title
        elif shape.placeholder_format and shape.placeholder_format.idx == 1:
            tf = shape.text_frame; tf.word_wrap = True
            while len(tf.paragraphs) > 1:
                tf.paragraphs[-1]._p.getparent().remove(tf.paragraphs[-1]._p)
            first_p = tf.paragraphs[0]; first_p.clear()
            first_p.alignment = PP_ALIGN.LEFT; first_p.level = lines[0][3]
            _add_paragraph_runs(first_p, lines[0][0], lines[0][1], lines[0][2])
            for line in lines[1:]:
                new_p = tf.add_paragraph(); new_p.alignment = PP_ALIGN.LEFT; new_p.level = line[3]
                _add_paragraph_runs(new_p, line[0], line[1], line[2])


def build_week12_presentation():
    """Build and return Week 12 as a standalone Presentation with 11 slides."""
    prs = Presentation(TEMPLATE)
    rid = prs.slides._sldIdLst[0].get("{http://schemas.openxmlformats.org/officeDocument/2006/relationships}id")
    prs.part.drop_rel(rid)
    del prs.slides._sldIdLst[0]

    s1 = prs.slides[0]
    for shape in s1.shapes:
        if shape.has_text_frame:
            for para in shape.text_frame.paragraphs:
                for run in para.runs:
                    if "Education Analyzer" in run.text:
                        run.text = "Housing Amenities Impact on Apartment\nListing Prices in the Southern USA"
                    elif "Week Ending" in run.text: run.text = f"Week Ending {WEEK_ENDING}"
                    elif "Submitted by" in run.text: run.text = f"Submitted by {SUBMITTED_BY}"
                    elif run.text.strip().startswith("May") and "Week" not in run.text: run.text = SUBMIT_DATE

    s2 = prs.slides[1]
    for shape in s2.shapes:
        if shape.has_text_frame and shape.placeholder_format:
            if shape.placeholder_format.idx == 0:
                for para in shape.text_frame.paragraphs:
                    for run in para.runs: run.text = "Agenda"
            elif shape.placeholder_format.idx == 1:
                for i, para in enumerate(shape.text_frame.paragraphs):
                    if i < len(AGENDA_ITEMS):
                        for run in para.runs: run.text = AGENDA_ITEMS[i]

    fill_content_slide(prs.slides[2], "Status Overview", STATUS_OVERVIEW_LINES)
    fill_content_slide(prs.slides[3], "Items Completed This Week", COMPLETED_LINES)
    fill_content_slide(prs.slides[4], "Items In Progress", IN_PROGRESS_LINES)
    fill_content_slide(prs.slides[5], "Items To Be Started", TBS_LINES)

    # Slide 7 - title text + screenshot image
    s7 = prs.slides[6]
    for shape in s7.shapes:
        if shape.has_text_frame and shape.placeholder_format:
            if shape.placeholder_format.idx == 0:
                for para in shape.text_frame.paragraphs:
                    for run in para.runs: run.text = "Sample of Items Completed This Week"
            elif shape.placeholder_format.idx == 1:
                tf = shape.text_frame; tf.word_wrap = True
                while len(tf.paragraphs) > 1:
                    tf.paragraphs[-1]._p.getparent().remove(tf.paragraphs[-1]._p)
                first_p = tf.paragraphs[0]; first_p.clear()
                first_p.alignment = PP_ALIGN.LEFT; first_p.level = 0
                r = first_p.add_run()
                r.text = "What: Final Presentation \u2014 Title Slide"
                r.font.size = Pt(11); r.font.bold = True; r.font.name = "Calibri"
                new_p = tf.add_paragraph(); new_p.alignment = PP_ALIGN.LEFT; new_p.level = 0
                r = new_p.add_run()
                r.text = "Why: Shows the completed presentation cover page; the full deliverable was submitted this week"
                r.font.size = Pt(10); r.font.bold = False; r.font.name = "Calibri"

    if os.path.exists(SCREENSHOT):
        slide_width = prs.slide_width
        img_left = Emu(599090)
        img_top = Emu(2300000)
        img_width = Emu(10993815)
        img_height = int(img_width * 0.5625)
        s7.shapes.add_picture(SCREENSHOT, img_left, img_top, img_width, img_height)

    fill_content_slide(prs.slides[7], "Issues, Risks, Concerns", IRC_LINES)
    fill_content_slide(prs.slides[8], "Next Steps", NEXT_STEPS_LINES)

    s10 = prs.slides[9]
    for shape in s10.shapes:
        if shape.has_text_frame and shape.placeholder_format:
            if shape.placeholder_format.idx == 0:
                for para in shape.text_frame.paragraphs:
                    for run in para.runs: run.text = "Personal Reflection"
            elif shape.placeholder_format.idx == 1:
                tf = shape.text_frame; tf.word_wrap = True
                while len(tf.paragraphs) > 1:
                    tf.paragraphs[-1]._p.getparent().remove(tf.paragraphs[-1]._p)
                first_p = tf.paragraphs[0]; first_p.clear()
                first_p.alignment = PP_ALIGN.LEFT; first_p.level = 0
                r = first_p.add_run(); r.text = REFLECTION
                r.font.size = Pt(11); r.font.bold = False; r.font.name = "Calibri"

    return prs


def merge_presentations(base_prs, append_path):
    """Merge slides from append_path PPTX after base_prs slides."""
    tmpdir = tempfile.mkdtemp()
    base_dir = os.path.join(tmpdir, "base")
    append_dir = os.path.join(tmpdir, "append")
    os.makedirs(base_dir); os.makedirs(append_dir)

    base_path = os.path.join(tmpdir, "base.pptx")
    base_prs.save(base_path)
    with zipfile.ZipFile(base_path, "r") as z: z.extractall(base_dir)
    with zipfile.ZipFile(append_path, "r") as z: z.extractall(append_dir)

    base_slides = sorted(
        int(re.match(r"slide(\d+)\.xml", f).group(1))
        for f in os.listdir(os.path.join(base_dir, "ppt/slides"))
        if f.endswith(".xml")
    )
    max_base_slide = max(base_slides) if base_slides else 0

    pres_path = os.path.join(base_dir, "ppt/presentation.xml")
    pres_tree = etree.parse(pres_path)
    sldIdLst = pres_tree.find(f"{{{NS['p']}}}sldIdLst")
    existing_ids = [int(s.get("id")) for s in sldIdLst]

    pres_rels_path = os.path.join(base_dir, "ppt/_rels/presentation.xml.rels")
    pres_rels_tree = etree.parse(pres_rels_path)
    max_rid = 0
    for rel in pres_rels_tree.getroot():
        m = re.match(r"rId(\d+)", rel.get("Id", ""))
        if m: max_rid = max(max_rid, int(m.group(1)))

    append_pres_path = os.path.join(append_dir, "ppt/presentation.xml")
    append_tree = etree.parse(append_pres_path)
    append_sldIdLst = append_tree.find(f"{{{NS['p']}}}sldIdLst")

    append_rels_path = os.path.join(append_dir, "ppt/_rels/presentation.xml.rels")
    append_rels_tree = etree.parse(append_rels_path)
    append_rid_to_slide = {}
    for rel in append_rels_tree.getroot():
        target = rel.get("Target", "")
        m = re.search(r"slides/slide(\d+)\.xml", target)
        if m:
            append_rid_to_slide[rel.get("Id")] = int(m.group(1))

    ordered_append_slides = []
    for sldId in append_sldIdLst:
        rid = sldId.get(f"{{{NS['r']}}}id")
        if rid in append_rid_to_slide:
            ordered_append_slides.append(append_rid_to_slide[rid])

    slide_num_offset = max_base_slide
    rid_offset = max_rid

    for i, old_num in enumerate(ordered_append_slides):
        new_num = slide_num_offset + 1 + i

        src = os.path.join(append_dir, f"ppt/slides/slide{old_num}.xml")
        dst = os.path.join(base_dir, f"ppt/slides/slide{new_num}.xml")
        shutil.copy2(src, dst)

        src_rels = os.path.join(append_dir, f"ppt/slides/_rels/slide{old_num}.xml.rels")
        dst_rels = os.path.join(base_dir, f"ppt/slides/_rels/slide{new_num}.xml.rels")
        if os.path.exists(src_rels):
            shutil.copy2(src_rels, dst_rels)

        p_sldId = etree.SubElement(sldIdLst, f"{{{NS['p']}}}sldId")
        p_sldId.set("id", str(max(existing_ids) + 1 + i))
        new_rid = f"rId{rid_offset + 1 + i}"
        p_sldId.set(f"{{{NS['r']}}}id", new_rid)

        rel_elem = etree.SubElement(pres_rels_tree.getroot(), "Relationship")
        rel_elem.set("Id", new_rid)
        rel_elem.set("Type", "http://schemas.openxmlformats.org/officeDocument/2006/relationships/slide")
        rel_elem.set("Target", f"slides/slide{new_num}.xml")

    pres_tree.write(pres_path, xml_declaration=True, encoding="UTF-8", standalone=True)
    pres_rels_tree.write(pres_rels_path, xml_declaration=True, encoding="UTF-8", standalone=True)

    ct_path = os.path.join(base_dir, "[Content_Types].xml")
    with open(ct_path, "r") as f: ct = f.read()
    for i in range(len(ordered_append_slides)):
        new_num = slide_num_offset + 1 + i
        override = '<Override PartName="/ppt/slides/slide{}.xml" ContentType="application/vnd.openxmlformats-officedocument.presentationml.slide+xml"/>'.format(new_num)
        if override not in ct:
            ct = ct.replace("</Types>", "  {}\n</Types>".format(override))
    with open(ct_path, "w") as f: f.write(ct)

    append_media_dir = os.path.join(append_dir, "ppt/media")
    base_media_dir = os.path.join(base_dir, "ppt/media")
    if os.path.exists(append_media_dir):
        os.makedirs(base_media_dir, exist_ok=True)
        for fname in os.listdir(append_media_dir):
            dst = os.path.join(base_media_dir, fname)
            if not os.path.exists(dst):
                shutil.copy2(os.path.join(append_media_dir, fname), dst)

    merged_pptx = os.path.join(tmpdir, "merged.pptx")
    with zipfile.ZipFile(merged_pptx, "w", zipfile.ZIP_DEFLATED) as zout:
        for root, dirs, files in os.walk(base_dir):
            for fname in files:
                fpath = os.path.join(root, fname)
                arcname = os.path.relpath(fpath, base_dir)
                zout.write(fpath, arcname)

    if os.path.exists(OUTPUT): os.remove(OUTPUT)
    shutil.copy2(merged_pptx, OUTPUT)
    shutil.rmtree(tmpdir, ignore_errors=True)
    return OUTPUT


# ===== MAIN =====
prs = build_week12_presentation()
print(f"Week 12 built: {len(prs.slides)} slides")

result = merge_presentations(prs, PREVIOUS)
prs_out = Presentation(result)
print(f"Final merged PPTX: {len(prs_out.slides)} total slides")
print(f"Saved to {OUTPUT}")
